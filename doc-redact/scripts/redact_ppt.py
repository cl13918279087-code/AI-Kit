#!/usr/bin/env python3
"""
redact_ppt.py - PPT 演示文稿脱敏脚本
支持 .pptx / .ppt（含 .ppt 需先转为 .pptx）
依赖: python-pptx
安装: pip install python-pptx
"""

import sys
import re
import zipfile
import shutil
import tempfile
from pathlib import Path
from xml.etree import ElementTree as ET

# ---------------------------------------------------------------------------
# 脱敏规则（9类 + 银行Logo）
# ---------------------------------------------------------------------------
REPLACEMENTS = [
    (re.compile(r'[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}'), 'XXXXX@XXXXX'),
    (re.compile(
        r'[^\x00-\xFF]{2,6}(?:省|自治区|市)?[^\x00-\xFF]{0,10}'
        r'(?:市|区)?[^\x00-\xFF]{0,10}'
        r'(?:街|路|道|巷|弄|号|大道|大街|东路|西路|南路|北路)[^\x00-\xFF]{0,30}'
    ), 'XX省XX市XX区XXXX'),
    (re.compile(r'[1-9]\d{5}(?:19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]'),
     'XXXXXXXXXXXXXXXXXX'),
    (re.compile(r'\b(?:\d{16}|\d{17}|\d{18}|\d{19})\b'), 'XXXXXXXXXXXXXXXX'),
    (re.compile(
        r'\d{4}[-年](?:0[1-9]|1[0-2])[-月](?:0[1-9]|[12]\d|3[01])[日]?\s*'
        r'|(?:19|20)\d{2}年\d{1,2}月\d{1,2}日'
    ), 'YYYY/MM/DD'),
    (re.compile(r'\b1[3-9]\d{9}\b'), 'XXXXXXXXXXX'),
    (re.compile(r'0\d{2,3}[-\s]?\d{7,8}'), '0XX-XXXXXXXX'),
    (re.compile(
        r'(?:(?:中国|交通|招商|浦发|兴业|民生|华夏|平安|光大|广发|浙商|渤海|恒丰|'
        r'南京|宁波|杭州|深圳|上海|北京|广州|农业|建设|工商|中国)银行|'
        r'(?:农信社|信用社|农商银行|合作银行|人民银行))'
    ), '[某银行]'),
    (re.compile(r'[\u4e00-\u9fa5]{2,4}(?![a-zA-Z0-9\u4e00-\u9fa5])'), 'XXX'),
]


def apply_redactions(text: str) -> str:
    if not text:
        return text
    result = text
    for pattern, replacement in REPLACEMENTS:
        result = pattern.sub(replacement, result)
    return result


def redact_pptx(input_path: str, output_path: str) -> None:
    """处理 .pptx 文件"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(input_path)
    redact_count = 0

    # 1. 处理所有幻灯片中的形状文本
    for slide_idx, slide in enumerate(prs.slides):
        # 收集该页含"银行"的形状位置，用于 Logo 检测
        bank_shape_positions = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        original = run.text
                        redacted = apply_redactions(original)
                        if redacted != original:
                            run.text = redacted
                            redact_count += 1
                    # 检查段落是否含银行关键词
                    para_text = ''.join(r.text for r in para.runs)
                    if '银行' in para_text:
                        bank_shape_positions.append(shape.left, shape.top)

            # 2. 处理表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                original = run.text
                                redacted = apply_redactions(original)
                                if redacted != original:
                                    run.text = redacted
                                    redact_count += 1

        # 3. 处理演讲者备注
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            for para in notes_frame.paragraphs:
                for run in para.runs:
                    original = run.text
                    redacted = apply_redactions(original)
                    if redacted != original:
                        run.text = redacted
                        redact_count += 1

    # 4. 处理页眉页脚
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            original = run.text
                            redacted = apply_redactions(original)
                            if redacted != original:
                                run.text = redacted
                                redact_count += 1

    # 5. 直接编辑 XML，处理文本框/形状深层内容 + 银行Logo图片
    _redact_pptx_xml(input_path, output_path, redact_count)
    prs.save(output_path)


def _redact_pptx_xml(input_path: str, output_path: str, base_count: int) -> int:
    """
    直接编辑 pptx 内部 XML，处理：
    - 组合形状内的文本
    - SmartArt（略过，因为结构复杂）
    - 页眉页脚文本
    - 演讲者备注
    - 银行 Logo 图片（检测 media 目录中文件名含银行关键词的图片）
    """
    import os

    count = base_count
    tmp_dir = tempfile.mkdtemp()
    extract_dir = Path(tmp_dir) / 'pptx_extracted'
    extract_dir.mkdir()

    with zipfile.ZipFile(input_path, 'r') as z:
        z.extractall(extract_dir)

    # 处理所有 ppt/slides/*.xml
    slides_dir = extract_dir / 'ppt' / 'slides'
    if slides_dir.exists():
        for xml_file in slides_dir.glob('*.xml'):
            content = xml_file.read_text(encoding='utf-8')
            original = content
            for pattern, replacement in REPLACEMENTS:
                content = pattern.sub(replacement, content)
            if content != original:
                xml_file.write_text(content, encoding='utf-8')
                count += 1

    # 处理备注页
    notes_dir = extract_dir / 'ppt' / 'notesSlides'
    if notes_dir.exists():
        for xml_file in notes_dir.glob('*.xml'):
            content = xml_file.read_text(encoding='utf-8')
            original = content
            for pattern, replacement in REPLACEMENTS:
                content = pattern.sub(replacement, content)
            if content != original:
                xml_file.write_text(content, encoding='utf-8')
                count += 1

    # 处理页眉页脚
    for subdir in ['slideMasters', 'slideLayouts']:
        dir_path = extract_dir / 'ppt' / subdir
        if dir_path.exists():
            for xml_file in dir_path.glob('*.xml'):
                content = xml_file.read_text(encoding='utf-8')
                original = content
                for pattern, replacement in REPLACEMENTS:
                    content = pattern.sub(replacement, content)
                if content != original:
                    xml_file.write_text(content, encoding='utf-8')
                    count += 1

    # 处理媒体文件中的银行 Logo
    media_dir = extract_dir / 'ppt' / 'media'
    if media_dir.exists():
        for img_file in media_dir.iterdir():
            img_name = img_file.name.lower()
            if any(k in img_name for k in ['bank', 'logo', '银行']):
                _replace_image_with_black(img_file)
                count += 1
                print(f"  [银行Logo图片替换] {img_file.name} → 纯黑图")

    # 重新打包（python-pptx 写入后再用 zip 覆盖含 Logo 的媒体文件）
    tmp_out = Path(tmp_dir) / 'output.pptx'
    with zipfile.ZipFile(tmp_out, 'w', zipfile.ZIP_DEFLATED) as z:
        for file_path in extract_dir.rglob('*'):
            if file_path.is_file():
                arcname = str(file_path.relative_to(extract_dir))
                z.write(file_path, arcname)

    # 用修改后的 zip 内容覆盖 python-pptx 保存的文件
    shutil.copy(tmp_out, output_path)
    shutil.rmtree(tmp_dir, ignore_errors=True)
    return count


def _replace_image_with_black(image_path: Path) -> None:
    """将图片替换为纯黑图"""
    from PIL import Image

    try:
        img = Image.open(image_path)
        w, h = img.size
        black = Image.new('RGB', (max(w, 10), max(h, 10)), (0, 0, 0))
        black.save(image_path)
    except Exception as e:
        print(f"  [警告] 无法处理图片 {image_path}: {e}")


def redact_ppt_to_pptx(input_path: str, output_path: str) -> None:
    """将 .ppt 转换为 .pptx 后处理"""
    import subprocess
    stem = Path(input_path).stem
    tmp_pptx = str(Path(input_path).with_name(f"{stem}_converted.pptx"))
    result = subprocess.run(
        ['textutil', '-convert', 'pptx', '-output', tmp_pptx, input_path],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        print(f"[错误] .ppt 转换失败: {result.stderr}", file=sys.stderr)
        sys.exit(1)
    redact_pptx(tmp_pptx, output_path)
    Path(tmp_pptx).unlink(missing_ok=True)


def main():
    if len(sys.argv) < 2:
        print("用法: python3 redact_ppt.py <输入文件路径> [输出文件路径]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    if output_file is None:
        stem = Path(input_file).stem
        output_file = str(Path(input_file).with_name(f"{stem}_脱敏.pptx"))

    ext = Path(input_file).suffix.lower()
    if ext == '.ppt':
        redact_ppt_to_pptx(input_file, output_file)
    else:
        redact_pptx(input_file, output_file)


if __name__ == '__main__':
    main()
